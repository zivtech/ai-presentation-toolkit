"""
Command Line Interface for Presentation Toolkit

Provides entry points for:
- pptx-migrate: Migrate presentations to branded templates
- pptx-analyze: Analyze brand compliance
- pptx-extract: Extract content from presentations
- pptx-diagnose: Run template diagnostics
"""

import sys
import json
import shutil
import tempfile
import argparse
from pathlib import Path
from typing import Optional

from .config import load_config, BrandConfig
from .migrate import migrate_presentation, migrate_from_content, detect_and_parse
from .analyze import analyze_presentation, get_analysis_json
from .extract import extract_pptx_to_markdown


def migrate_command(args: argparse.Namespace) -> int:
    """Execute migrate command."""
    print("=" * 60)
    print("Presentation Migration")
    print("=" * 60)

    # Load configuration
    config = load_config(args.config)
    print(f"Brand: {config.brand_name}")

    # Determine template path
    template_path = args.template
    if template_path is None:
        if config.template.default:
            template_path = config.template.default
        else:
            print("Error: No template specified. Use --template or set template.default in config.")
            return 1

    try:
        from_content = getattr(args, 'from_content', None)

        if from_content:
            # Skip parsing — use pre-built content.json
            print(f"Loading content from: {from_content}")
            migrate_from_content(
                from_content,
                args.output,
                config,
                template_path,
                insert_images=not args.no_images,
                diagnose=True,
                use_cookbook=getattr(args, 'use_cookbook', False),
            )
            return 0

        # Create temp directory for extracted images
        image_dir = None
        if not args.no_images:
            image_dir = Path(tempfile.mkdtemp(prefix='migrate_images_'))
            print(f"Image extraction: enabled (temp: {image_dir})")
        else:
            print("Image extraction: disabled")

        # Parse input
        slides = detect_and_parse(args.input, image_output_dir=image_dir)
        print(f"Parsed {len(slides)} slides from input")

        # Count slides with images
        slides_with_images = sum(1 for s in slides if s.get('images'))
        if slides_with_images > 0:
            print(f"  {slides_with_images} slides have extractable images")

        # Optionally save intermediate content.json
        save_content = getattr(args, 'save_content', None)
        if save_content:
            from .content import slides_to_content_document, save_content_document
            source_file = str(Path(args.input).name)
            source_format = Path(args.input).suffix.lstrip('.')
            doc = slides_to_content_document(slides, source_file, source_format)
            save_content_document(doc, save_content)
            print(f"Saved intermediate content to: {save_content}")

        # Migrate
        migrate_presentation(
            slides,
            args.output,
            config,
            template_path,
            insert_images=not args.no_images,
            diagnose=True,
            use_cookbook=getattr(args, 'use_cookbook', False),
        )

        # Cleanup temp image directory
        if image_dir and image_dir.exists():
            shutil.rmtree(image_dir)

        return 0

    except Exception as e:
        print(f"\nError: {e}")
        if args.verbose:
            import traceback
            traceback.print_exc()
        return 1


def analyze_command(args: argparse.Namespace) -> int:
    """Execute analyze command."""
    print("=" * 60)
    print("Brand Compliance Analysis")
    print("=" * 60)

    config = load_config(args.config)
    print(f"Brand: {config.brand_name}")

    try:
        issues = analyze_presentation(args.input, config, verbose=not args.json)

        if args.json:
            # Get total slide count
            try:
                from pptx import Presentation
                prs = Presentation(args.input)
                total_slides = len(prs.slides)
            except:
                total_slides = len(issues)

            result = get_analysis_json(issues, total_slides)
            print(json.dumps(result, indent=2))

        if args.strict and issues:
            print(f"\nStrict mode: {len(issues)} compliance issues found")
            return 1

        return 0

    except Exception as e:
        print(f"\nError: {e}")
        if args.verbose:
            import traceback
            traceback.print_exc()
        return 1


def extract_command(args: argparse.Namespace) -> int:
    """Execute extract command."""
    print("=" * 60)
    print("Content Extraction")
    print("=" * 60)

    try:
        extract_pptx_to_markdown(
            args.input,
            args.output,
            extract_images=args.images
        )
        print("\nExtraction complete!")
        return 0

    except Exception as e:
        print(f"\nError: {e}")
        if args.verbose:
            import traceback
            traceback.print_exc()
        return 1


def diagnose_command(args: argparse.Namespace) -> int:
    """Execute diagnose command."""
    from .diagnose import diagnose_template

    config = None
    if args.config:
        config = load_config(args.config)

    try:
        report = diagnose_template(args.template, config=config)

        if getattr(args, 'json', False):
            print(json.dumps(report.to_dict(), indent=2))
        else:
            report.print_report()

        if args.strict and report.has_blocking_issues:
            return 1

        return 0

    except Exception as e:
        print(f"\nError: {e}")
        if getattr(args, 'verbose', False):
            import traceback
            traceback.print_exc()
        return 1


def main():
    """Main CLI entry point."""
    parser = argparse.ArgumentParser(
        description='Presentation Toolkit - Brand-aware presentation migration and analysis',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s migrate input.pptx output.pptx --config brand.yaml
  %(prog)s migrate --from-content content.json output.pptx --config brand.yaml
  %(prog)s analyze deck.pptx --config brand.yaml --json
  %(prog)s extract input.pptx --output content.md --images
  %(prog)s diagnose template.pptx --config brand.yaml
        """
    )

    parser.add_argument('-v', '--verbose', action='store_true', help='Show detailed output')

    subparsers = parser.add_subparsers(dest='command', help='Available commands')

    # Migrate command
    migrate_parser = subparsers.add_parser('migrate', help='Migrate presentation to branded template')
    migrate_parser.add_argument('input', nargs='?', help='Input file (PPTX, PDF, MD, or CSV)')
    migrate_parser.add_argument('output', help='Output PPTX file')
    migrate_parser.add_argument('--config', '-c', required=True, help='Brand configuration file (YAML/JSON)')
    migrate_parser.add_argument('--template', '-t', help='Template PPTX (overrides config)')
    migrate_parser.add_argument('--no-images', action='store_true', help='Skip image extraction/insertion')
    migrate_parser.add_argument('--save-content', metavar='PATH', help='Save intermediate content.json before migration')
    migrate_parser.add_argument('--from-content', metavar='PATH', help='Skip parsing, use pre-built content.json')
    migrate_parser.add_argument('--use-cookbook', action='store_true', help='Force cookbook mode (absolute positioning, no template placeholders)')
    migrate_parser.add_argument('-v', '--verbose', action='store_true', help='Show detailed output')

    # Analyze command
    analyze_parser = subparsers.add_parser('analyze', help='Analyze presentation for brand compliance')
    analyze_parser.add_argument('input', help='PPTX file to analyze')
    analyze_parser.add_argument('--config', '-c', required=True, help='Brand configuration file (YAML/JSON)')
    analyze_parser.add_argument('--json', action='store_true', help='Output results as JSON')
    analyze_parser.add_argument('--strict', action='store_true', help='Exit with error if issues found')
    analyze_parser.add_argument('-v', '--verbose', action='store_true', help='Show detailed output')

    # Extract command
    extract_parser = subparsers.add_parser('extract', help='Extract content from presentation')
    extract_parser.add_argument('input', help='PPTX file to extract from')
    extract_parser.add_argument('--output', '-o', help='Output markdown file (default: input.md)')
    extract_parser.add_argument('--images', action='store_true', help='Also extract images')
    extract_parser.add_argument('-v', '--verbose', action='store_true', help='Show detailed output')

    # Diagnose command
    diagnose_parser = subparsers.add_parser('diagnose', help='Run template diagnostics')
    diagnose_parser.add_argument('template', help='Template PPTX file to diagnose')
    diagnose_parser.add_argument('--config', '-c', help='Brand configuration file (YAML/JSON)')
    diagnose_parser.add_argument('--strict', action='store_true', help='Exit with error if blocking issues found')
    diagnose_parser.add_argument('--json', action='store_true', help='Output results as JSON')
    diagnose_parser.add_argument('-v', '--verbose', action='store_true', help='Show detailed output')

    args = parser.parse_args()

    if args.command is None:
        parser.print_help()
        return 1

    if args.command == 'migrate':
        return migrate_command(args)
    elif args.command == 'analyze':
        return analyze_command(args)
    elif args.command == 'extract':
        return extract_command(args)
    elif args.command == 'diagnose':
        return diagnose_command(args)
    else:
        parser.print_help()
        return 1


# Entry points for direct script execution
def pptx_migrate():
    """Entry point for pptx-migrate command."""
    sys.argv = ['pptx-migrate', 'migrate'] + sys.argv[1:]
    return main()


def pptx_analyze():
    """Entry point for pptx-analyze command."""
    sys.argv = ['pptx-analyze', 'analyze'] + sys.argv[1:]
    return main()


def pptx_extract():
    """Entry point for pptx-extract command."""
    sys.argv = ['pptx-extract', 'extract'] + sys.argv[1:]
    return main()


def pptx_diagnose():
    """Entry point for pptx-diagnose command."""
    sys.argv = ['pptx-diagnose', 'diagnose'] + sys.argv[1:]
    return main()


if __name__ == '__main__':
    sys.exit(main())
